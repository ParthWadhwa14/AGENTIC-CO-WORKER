from pathlib import Path
from uuid import uuid4

from pptx import Presentation

from app.models import ParsedElement


class PPTXLoader:
    source_type = "pptx"

    def load(self, file_path: str, document_id: str | None = None) -> list[ParsedElement]:
        document_id = document_id or str(uuid4())
        path = Path(file_path)

        presentation = Presentation(file_path)
        elements: list[ParsedElement] = []

        slide_titles = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_text_parts = []
            title = None

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                text = shape.text.strip()

                if not text:
                    continue

                if title is None:
                    title = text.splitlines()[0].strip()

                slide_text_parts.append(text)

            if title:
                slide_titles.append(f"Slide {slide_index}: {title}")

            slide_text = "\n".join(slide_text_parts).strip()

            if not slide_text:
                continue

            formatted_text = f"""
Slide {slide_index}
Title: {title or "Untitled"}

Content:
{slide_text}
""".strip()

            elements.append(
                ParsedElement(
                    document_id=document_id,
                    source_type=self.source_type,
                    file_name=path.name,
                    slide_number=slide_index,
                    element_type="slide",
                    text=formatted_text,
                    metadata={
                        "file_path": str(path),
                        "slide_title": title,
                    },
                )
            )

        # Deck summary chunk
        if slide_titles:
            elements.insert(
                0,
                ParsedElement(
                    document_id=document_id,
                    source_type=self.source_type,
                    file_name=path.name,
                    element_type="deck_summary",
                    text=(
                        f"PowerPoint deck: {path.name}\n"
                        f"Number of slides: {len(presentation.slides)}\n\n"
                        f"Slide titles:\n" + "\n".join(slide_titles)
                    ),
                    metadata={
                        "file_path": str(path),
                        "slide_count": len(presentation.slides),
                    },
                )
            )

        return elements